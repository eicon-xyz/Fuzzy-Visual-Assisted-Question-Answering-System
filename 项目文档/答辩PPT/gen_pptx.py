# -*- coding: utf-8 -*-
"""生成 HAJIMI 验收答辩 PPTX（16:9），嵌入 images/ 图表，为运行截图预留占位框。"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(BASE, "images")

# 配色
BG = RGBColor(0x0F, 0x17, 0x2A)
CARD = RGBColor(0x1E, 0x29, 0x3B)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
SLATE = RGBColor(0x64, 0x74, 0x8B)
FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def txt(slide, x, y, w, h, lines, size=18, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, space=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space is not None:
            p.space_after = Pt(space)
        # 支持 (text, color, size, bold) 元组
        if isinstance(ln, tuple):
            t, c, s, b = (list(ln) + [color, size, bold])[:4]
        else:
            t, c, s, b = ln, color, size, bold
        r = p.add_run(); r.text = t
        r.font.size = Pt(s); r.font.bold = b
        r.font.color.rgb = c; r.font.name = font
    return tb


def bar(slide, color=CYAN, y=1.15, h=0.06, x=0.7, w=3.0):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background()
    return sp


def title(slide, text, kicker=None):
    if kicker:
        txt(slide, 0.7, 0.32, 8, 0.4, kicker, size=13, color=CYAN, bold=True)
    txt(slide, 0.7, 0.55, 12, 0.8, text, size=30, color=WHITE, bold=True)
    bar(slide)


def pageno(slide, n):
    txt(slide, 12.4, 7.02, 0.8, 0.4, f"{n:02d}/13", size=11, color=SLATE, align=PP_ALIGN.RIGHT)


def pic(slide, name, x, y, w, h=None):
    path = os.path.join(IMG, name)
    if os.path.exists(path):
        kw = dict(width=Inches(w))
        if h:
            kw["height"] = Inches(h)
        slide.shapes.add_picture(path, Inches(x), Inches(y), **kw)
        return True
    return False


def placeholder(slide, x, y, w, h, label):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = CARD
    sp.line.color.rgb = SLATE; sp.line.width = Pt(1.25); sp.line.dash_style = 2
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "📷 " + label
    r.font.size = Pt(15); r.font.color.rgb = SLATE; r.font.name = FONT
    return sp


def bullets(slide, x, y, w, h, items, size=17, gap=8, color=WHITE):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        rb = p.add_run(); rb.text = "▪ "; rb.font.size = Pt(size)
        rb.font.color.rgb = CYAN; rb.font.name = FONT; rb.font.bold = True
        r = p.add_run(); r.text = it; r.font.size = Pt(size)
        r.font.color.rgb = color; r.font.name = FONT
    return tb


# ---------- P1 封面 ----------
s = prs.slides.add_slide(BLANK); bg(s)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), SW, Inches(2.6))
band.fill.solid(); band.fill.fore_color.rgb = CARD; band.line.fill.background()
txt(s, 0.8, 2.55, 11.7, 1.2, "HAJIMI 自动操作助手", size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, 0.8, 3.85, 11.7, 0.6, "看得懂屏幕 · 听得懂指令 · 帮你动手做", size=22, color=CYAN, align=PP_ALIGN.CENTER)
txt(s, 0.8, 4.55, 11.7, 0.5, "屏幕感知  →  意图理解  →  安全执行", size=16, color=SLATE, align=PP_ALIGN.CENTER)
txt(s, 0.8, 6.4, 11.7, 0.6,
    "32 组  ·  潘振喆 / 杨名 / 涂浚稷  ·  2026.07", size=15, color=WHITE, align=PP_ALIGN.CENTER)
placeholder(s, 4.9, 0.5, 3.5, 1.8, "IMG-P01 项目代表图")

# ---------- P2 背景痛点 ----------
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "项目背景与痛点", "WHY · 为什么做")
pic(s, "P02_pain_vs_hajimi.png", 1.4, 1.7, 10.5)
txt(s, 0.7, 6.7, 12, 0.5,
    "面向新手 / 老年 / 视障用户：把「查教程—切窗口—反复试」压缩成「说一句话」",
    size=15, color=SLATE, align=PP_ALIGN.CENTER)
pageno(s, 2)

# ---------- P3 定位能力 ----------
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "系统定位与核心能力", "WHAT · 做什么")
pic(s, "P03_capability_chain.png", 1.4, 1.7, 10.5)
pageno(s, 3)

# ---------- P4 架构 ----------
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "技术架构 · 四层逻辑", "HOW · 技术实现")
pic(s, "P04_architecture.png", 1.7, 1.55, 9.9)
pageno(s, 4)

# ---------- P5 流程 ----------
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "核心业务流程 · 时序", "HOW · 业务链路")
pic(s, "P05_sequence.png", 1.7, 1.55, 9.9)
pageno(s, 5)

# ---------- P6 A端 ----------
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "A 端 · 后端服务", "成果 · FastAPI")
pic(s, "P06_endpoints.png", 1.9, 1.55, 9.5)
pageno(s, 6)

# ---------- P7 B端 ----------
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "B 端 · 桌面客户端", "成果 · PyQt5")
bullets(s, 0.75, 1.75, 5.6, 5, [
    "无边框悬浮面板 480×520 + 紧凑条 320×52",
    "三区：输入栏 / 步骤卡片 / 截图+日志",
    "步骤 6 态：待执行/执行中/完成/失败/拦截/跳过",
    "透明覆盖层：红色箭头 + 高亮框 + 编号标签",
    "多主题 · 系统托盘 · 拖拽 · 置顶 · Resize",
], size=17, gap=12)
placeholder(s, 6.7, 1.7, 6.0, 5.1, "IMG-P07 悬浮面板 + 全屏覆盖层截图")
pageno(s, 7)

# ---------- P8 自动执行 ----------
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "B 端 · V2.2 自动执行", "★ 本期亮点")
bullets(s, 0.75, 1.75, 5.6, 5, [
    "V1 只「指路」→ V2.2 直接代为操作 (pyautogui)",
    "每步卡片显示 action 类型 + 风险评分圆点",
    "L2 快路径 <3s · L3 慢路径 LLM 规划 5–10s",
    "蓝图状态机：生成→执行→完成/挂起/回退",
    "风险 ≥4 自动挂起，等待用户确认",
], size=17, gap=12)
placeholder(s, 6.7, 1.7, 6.0, 5.1, "IMG-P08 执行状态卡片 + 风险评分 + 执行日志")
pageno(s, 8)

# ---------- P9 C端 ----------
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "C 端 · Web 管理面板", "成果 · Vue3 + ECharts")
bullets(s, 0.75, 1.75, 5.6, 5, [
    "Vue3 + Vite + Element Plus + ECharts + Pinia",
    "6 页：登录/总览/失败归因/数据流/配置/健康",
    "总览：5 KPI + 反馈饼图 + L2/L3 + 24h 趋势",
    "失败归因下钻单任务 + LLM 快照",
    "配置热部署 · 健康告警 · CSV 导出",
], size=17, gap=12)
placeholder(s, 6.7, 1.7, 6.0, 5.1, "IMG-P09 总览 / 失败归因 / 数据流 三页拼图")
pageno(s, 9)

# ---------- P10 数据连通 ----------
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "端到端数据连通性验证", "成果 · 全链路")
placeholder(s, 0.75, 1.75, 6.4, 5.0, "IMG-P10 connectivity 终端输出 (全 PASS)")
bullets(s, 7.5, 1.9, 5.2, 5, [
    "审计回路：C→POST→A→DB→GET→C",
    "配置回路：C→deploy→A→DB→pull→C",
    "B-C 链路：9 信号契约仿真联调",
    "鉴权：无 Key→401 · JWT 正常签发",
    "脚本化：一键复验 real_api_test.py",
], size=16, gap=13)
pageno(s, 10)

# ---------- P11 测试 ----------
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "测试与质量保障", "成果 · 237 项")
pic(s, "P11_test_matrix.png", 1.7, 1.6, 9.9)
pageno(s, 11)

# ---------- P12 总结 ----------
s = prs.slides.add_slide(BLANK); bg(s)
title(s, "项目总结与展望", "总结")
pic(s, "P12_summary.png", 1.4, 1.7, 10.5)
pageno(s, 12)

# ---------- P13 致谢 ----------
s = prs.slides.add_slide(BLANK); bg(s)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.2), SW, Inches(1.5))
band.fill.solid(); band.fill.fore_color.rgb = CARD; band.line.fill.background()
txt(s, 0.8, 2.35, 11.7, 1.1, "谢谢！欢迎提问", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, 0.8, 4.2, 11.7, 0.5, "团队分工", size=18, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
txt(s, 0.8, 4.75, 11.7, 1.6, [
    "潘振喆（20230383）— A 端 后端 / AI 规划",
    "杨名（20230645）— B 端 前端 / 桌面客户端",
    "涂浚稷（20230353）— C 端 集成 / 语音 / 管理面板",
], size=17, color=WHITE, align=PP_ALIGN.CENTER, space=6)

out = os.path.join(BASE, "HAJIMI-答辩PPT.pptx")
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
