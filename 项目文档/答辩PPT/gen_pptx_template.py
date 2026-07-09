# -*- coding: utf-8 -*-
"""
HAJIMI 答辩 PPT —— 严格沿用《PPTer吧-毕业论文答辩PPT免费模板.pptx》模板文件格式
在模板基础上：
  1) 逐页替换占位文字为 HAJIMI 讲解性内容，并按框尺寸设定字号 + 自动换行 + 缩放，杜绝重叠；
  2) 适当加多文字描述；
  3) 演示章节插入多图页（每页 2–5 张），不再一页一张。
输出：项目文档/答辩PPT/HAJIMI-答辩PPT-模板版.pptx（若被占用则存 _new）
"""
import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

BASE = os.path.dirname(os.path.abspath(__file__))
TPL = os.path.join(BASE, "..", "模版", "PPTer吧-毕业论文答辩PPT免费模板.pptx")
SHOT = os.path.join(BASE, "..", "实训项目图片")
OUT = os.path.join(BASE, "HAJIMI-答辩PPT-模板版.pptx")

NAVY = RGBColor(0x20, 0x51, 0x7C)
GOLD = RGBColor(0xFF, 0xC0, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUB = RGBColor(0x60, 0x70, 0x80)
FONT = "微软雅黑"
SW, SH = Inches(13.333), Inches(7.5)

# 字号约定
TITLE = 22   # 内容页节标题
SUBH = 14    # 小标题
BODY = 13    # 正文
BUB = 11     # 气泡短句
KW = 14      # 关键词

# ── 文字替换：{slide: {shape: 文本 或 (文本,字号)}} ──
REPL = {
    0: {0: "HAJIMI 自动操作助手", 1: "项目：模糊视觉辅助问答系统", 2: "团队：第 32 组",
        3: "答辩人：潘振喆 / 杨名 / 涂浚稷", 4: "汇报日期：2026 年 7 月"},
    1: {12: "项目背景与意义", 13: "系统介绍与核心能力", 14: "功能演示 · 桌面端",
        15: "管理与运营数据", 16: "项目成果与亮点", 17: "总结与展望"},
    2: {2: "项目背景与意义"},
    3: {1: ("用户面临的困境", TITLE),
        4: ("操作门槛高", SUBH),
        5: ("面对复杂的软件界面，用户常常“知道要做什么，却找不到在哪里点”。功能层层嵌套、"
            "按钮命名专业，很多操作要在多级菜单里反复翻找；即便是简单任务也变得困难，还容易点错。", BODY),
        6: ("求助方式低效", SUBH),
        7: ("查帮助文档、按 F1、上网搜教程、在多个窗口之间来回切换——不仅费时费力，网上的教程还"
            "常常与自己电脑上的软件版本、界面布局对不上号，照着做也走不通。", BODY)},
    4: {1: ("谁最需要它", TITLE),
        7: ("HAJIMI 面向几类典型用户，帮助他们跨过操作门槛，把“不会用”变成“说一句就行”。", BODY),
        8: ("新手用户：初次接触软件，面对陌生界面无从下手，容易在复杂菜单里迷路。", BUB),
        9: ("老年用户：视力与精细操作能力下降，看不清、点不准，需要更直观的引导。", BUB),
        10: ("视障用户：读屏工具对图形界面覆盖有限，很多按钮难以感知与操作。", BUB),
        11: ("忙碌用户：不想研究操作细节，只想一句话交代目标，剩下的交给系统。", BUB)},
    5: {1: ("我们想做什么", TITLE),
        9: ("说一句话就行：用最自然的语言表达目标，无需记忆菜单路径。", BUB),
        10: ("系统自己看懂：自动识别屏幕元素、理解意图并规划操作步骤。", BUB),
        14: ("在界面上帮忙：既能标注指路，也能直接代替用户点击、输入。", BUB),
        15: ("始终安全可控：风险评分与红线拦截贯穿全程，高危操作需确认。", BUB)},
    6: {2: "系统介绍与核心能力"},
    7: {1: ("三步走，六件事", TITLE),
        2: ("① 屏幕感知：截取屏幕并识别可点、可输入的界面元素。", BUB),
        3: ("② 意图理解：把自然语言归类到意图域并消解模糊指代。", BUB),
        4: ("③ 步骤规划：由多模态大模型生成带坐标的操作步骤。", BUB),
        5: ("④ 可视标注：在真实界面叠加箭头、高亮框与编号提示。", BUB),
        6: ("⑤ 自动执行：调用系统自动化能力代为点击、输入、双击。", BUB),
        7: ("⑥ 安全护栏：风险评分与红线拦截双重把关，高危自动挂起。", BUB)},
    8: {1: ("两种使用方式，都很自然", TITLE),
        7: ("语音说", KW), 8: ("打字问", KW), 9: ("看标注", KW), 10: ("帮你做", KW),
        11: ("你既可以按住麦克风用语音说出目标，也可以直接在输入框打字。系统会先弄清你要做什么，"
             "再决定怎么帮：简单任务在界面上标注“点这里”，复杂任务则一步步代你完成。", BODY),
        12: ("整个过程可暂停、可恢复、可随时停止，主动权始终在用户手中。", BODY)},
    9: {2: "功能演示 · 桌面端"},
    10: {1: ("桌面助手：贴身又不打扰", TITLE),
         6: ("常驻悬浮面板", SUBH), 7: ("停靠在桌面角落，随时输入指令、查看每一步的执行进度与状态。", 12),
         8: ("紧凑浮动条", SUBH), 9: ("一键收起为一条小窗，仅保留输入与状态，完全不遮挡工作区。", 12),
         14: ("多套界面主题", SUBH), 15: ("内置默认蓝、典雅黑、牛皮纸、黑金轻奢与适老等多套主题。", 12),
         16: ("屏幕标注覆盖层", SUBH), 17: ("把“下一步点这里”用箭头、高亮框与编号直接画在真实界面上。", 12)},
    13: {1: "PART  FOUR", 2: "管理与运营数据"},
    14: {1: ("Web 管理面板：一屏掌握全局", TITLE),
         3: ("总览仪表盘：任务量、成功率、耗时与趋势一目了然。", 12),
         4: ("失败归因：按类型下钻到每个失败任务，还原当时的 LLM 快照。", 12),
         5: ("数据流监控：以桑基图与双轴曲线呈现链路流向与调用量。", 12),
         6: ("系统配置：各类运行参数在线热部署，保存即生效并留痕。", 12),
         7: ("用户管理：用户、角色、任务数与登录情况集中维护。", 12),
         8: ("健康监控：组件状态、资源占用与告警实时掌握并可导出。", 12),
         9: ("六大页面协同，覆盖运营分析与系统运维的日常所需。", 12)},
    16: {2: "项目成果与亮点"},
    17: {1: ("五大亮点", TITLE),
         8: ("看得懂屏幕：多模态视觉理解真实界面，不依赖写死坐标。", BUB),
         9: ("帮得上忙：从“指路”升级到“代做”，真正替用户完成操作。", BUB),
         10: ("用得放心：风险评分（1–5）叠加红线拦截，双重护栏。", BUB),
         11: ("随手可用：常驻桌面悬浮，语音与文字输入皆可。", BUB),
         12: ("可管可控：Web 面板对任务、失败与配置全程可视化运营。", BUB)},
    18: {1: ("三端联通，稳定可用", TITLE),
         3: ("桌面客户端：悬浮面板、标注覆盖、自动执行与多套主题均已完成并可稳定运行。", BODY),
         4: ("服务端：任务规划、安全校验、审计与配置等 31 个接口全部就位并通过联调。", BODY),
         5: ("管理面板：六大页面全部打通，数据从采集到展示端到端可追溯。", BODY)},
    19: {1: ("质量有保障", TITLE),
         3: ("系统经过完整的功能测试与全链路联调，既覆盖正常操作流程，也覆盖断网、误操作、"
             "被安全拦截等异常场景，端到端的数据连通性均已验证通过。", BODY),
         4: ("237 项测试全部通过", 13), 5: ("六大页面回归通过", 13), 6: ("关键链路端到端打通", 13)},
    20: {2: "总结与展望"},
    21: {1: ("一句话总结", TITLE),
         3: ("HAJIMI 让普通用户用最自然的方式——说一句话——就能完成原本繁琐的桌面操作，"
             "把“会不会用软件”这道门槛降到最低。", BODY),
         4: ("它看得懂屏幕、听得懂指令、帮得上忙，并且始终在安全边界内工作，让自动化既好用又可控。", BODY)},
    22: {1: ("未来展望", TITLE),
         4: ("更聪明", SUBH),
         5: ("持续打磨语音识别与合成的全链路体验，并提升复杂任务的规划成功率，让理解与执行更稳更准。", BODY),
         6: ("更普适", SUBH),
         7: ("适配多屏与高分辨率显示环境，扩展对更多软件与业务场景的支持，覆盖更广的用户群体。", BODY)},
    24: {0: "感谢聆听 · 敬请指正", 4: "第 32 组", 5: "潘振喆 / 杨名 / 涂浚稷", 8: "HAJIMI 自动操作助手"},
}
DROP = {11, 12, 15, 23}

# 多图页：(PART, 标题, [(图, 说明)...])
GALLERIES = [
    ("03", "五套界面主题", [("默认蓝主页面.png", "默认蓝"), ("典雅黑主页面.png", "典雅黑"),
        ("牛皮纸主页面.png", "牛皮纸"), ("黑金轻奢主页面.png", "黑金轻奢"), ("耄耋主页面.png", "耄耋（适老）")]),
    ("03", "系统设置 · 灵活可配", [("系统设置-主题默认蓝.png", "主题外观"),
        ("系统设置-语音.png", "语音交互"), ("系统设置-部署模式.png", "部署模式")]),
    ("03", "自动执行与安全护栏", [("L4检测确认.png", "执行前检测确认"),
        ("L5演示.png", "全屏标注 · 自动演示"), ("L5警示.png", "高风险安全警示")]),
    ("04", "登录与运营总览", [("web端登录页.png", "登录"), ("web端总览页.png", "总览仪表盘")]),
    ("04", "失败归因：从列表到详情", [("web端失败归因页.png", "失败列表与图表"),
        ("web端失败归因页详情.png", "详情抽屉（含 LLM 快照）")]),
    ("04", "数据流与用户管理", [("web端数据流监控.png", "数据流监控"), ("web端用户管理页.png", "用户管理")]),
    ("04", "系统配置与健康监控", [("web端系统配置页1.png", "系统配置"), ("web端健康监控页.png", "健康监控")]),
]


def set_text(shape, val):
    if not shape.has_text_frame:
        return
    text, size = (val if isinstance(val, tuple) else (val, None))
    tf = shape.text_frame
    tf.word_wrap = True
    for extra in list(tf.paragraphs[1:]):
        extra._p.getparent().remove(extra._p)
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = text
        for r in p0.runs[1:]:
            r.text = ""
    else:
        p0.add_run().text = text
    if size:
        for r in p0.runs:
            r.font.size = Pt(size)
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # 溢出时自动缩字，防重叠
    except Exception:
        pass


def add_gallery(prs, blank_idx, part, title, items):
    s = prs.slides.add_slide(prs.slide_layouts[blank_idx])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background(); bar.shadow.inherit = False
    lab = s.shapes.add_textbox(Inches(0.55), Inches(0.16), Inches(2), Inches(0.4))
    r = lab.text_frame.paragraphs[0].add_run(); r.text = "PART " + part
    r.font.size = Pt(12.5); r.font.color.rgb = GOLD; r.font.bold = True; r.font.name = FONT
    tb = s.shapes.add_textbox(Inches(0.55), Inches(0.48), Inches(11.5), Inches(0.6))
    rr = tb.text_frame.paragraphs[0].add_run(); rr.text = title
    rr.font.size = Pt(24); rr.font.color.rgb = WHITE; rr.font.bold = True; rr.font.name = FONT
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.03), Inches(2.6), Pt(4))
    ln.fill.solid(); ln.fill.fore_color.rgb = GOLD; ln.line.fill.background(); ln.shadow.inherit = False
    # 图片行
    ax, ay, aw, ah, cap = 0.5, 1.55, 12.33, 5.4, 0.5
    n = len(items); gap = 0.35
    cw = (aw - (n - 1) * gap) / n
    for i, (img, capt) in enumerate(items):
        p = os.path.join(SHOT, img)
        iw, ih = Image.open(p).size
        cell_x = ax + i * (cw + gap)
        scale = min(cw / (iw / 96), (ah - cap) / (ih / 96))
        fw, fh = (iw / 96) * scale, (ih / 96) * scale
        s.shapes.add_picture(p, Inches(cell_x + (cw - fw) / 2), Inches(ay + (ah - cap - fh) / 2),
                             width=Inches(fw), height=Inches(fh))
        cb = s.shapes.add_textbox(Inches(cell_x), Inches(ay + ah - cap + 0.05), Inches(cw), Inches(cap))
        cp = cb.text_frame.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = capt
        cr.font.size = Pt(12.5); cr.font.color.rgb = SUB; cr.font.bold = True; cr.font.name = FONT
    return s


def main():
    shutil.copyfile(TPL, OUT if not os.path.exists(OUT) else OUT)  # 覆盖
    prs = Presentation(TPL)
    # 1) 文字替换 + 字号/换行/缩放
    for idx, mapping in REPL.items():
        shapes = list(prs.slides[idx].shapes)
        for si, val in mapping.items():
            if si < len(shapes):
                set_text(shapes[si], val)
    # 2) 找空白版式
    blank_idx = 5
    for i, l in enumerate(prs.slide_layouts):
        if len(l.placeholders) == 0 and len(l.shapes) == 0:
            blank_idx = i; break
    # 3) 追加多图页
    for part, title, items in GALLERIES:
        add_gallery(prs, blank_idx, part, title, items)
    # 4) 重排 + 丢弃不合适页
    ids = list(prs.slides._sldIdLst)  # 0..24 模板 + 25..31 多图
    g = list(range(25, 25 + len(GALLERIES)))
    t3, t4 = g[:3], g[3:]
    desired = ([0, 1, 2, 3, 4, 5, 6, 7, 8, 9, 10] + t3 + [13, 14] + t4
               + [16, 17, 18, 19, 20, 21, 22, 24])
    sldIdLst = prs.slides._sldIdLst
    for c in list(sldIdLst):
        sldIdLst.remove(c)
    for i in desired:
        sldIdLst.append(ids[i])
    # 5) 保存（处理占用）
    out = OUT
    try:
        prs.save(out)
    except PermissionError:
        out = OUT.replace(".pptx", "_new.pptx")
        prs.save(out)
        print("原文件被占用，已存为:", os.path.basename(out))
    print("saved", os.path.basename(out), "| slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
